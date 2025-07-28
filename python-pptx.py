from pptx import Presentation
from pptx.util import Inches

# Créer une présentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Fonction pour ajouter une diapositive
def add_slide(title, content):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Diapositive de titre
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Cahier des Charges - EventMagic"
slide.placeholders[1].text = "Application de Réservation d'Événements"

# Slides de contenu
sections = [
    ("1. Présentation du Projet", 
     "EventMagic est une plateforme web pour organiser des événements privés.\n\nObjectif principal : Centraliser l'organisation d'événements privés.\nObjectifs secondaires :\n- Réservation de professionnels\n- Gestion des invitations\n- Hébergements\n- Cagnotte\n\nCible :\n- Particuliers (25-55 ans)\n- Professionnels\n- Invités"),
    
    ("2.1 Fonctionnalités Principales", 
     "• Gestion des événements : type, date, lieu, budget, timeline\n• Catalogue de professionnels : fiches, filtres, géolocalisation\n• Réservation : devis, contrat, paiement sécurisé\n• Invitations : envoi, suivi, relance, régimes alimentaires\n• Hébergements : APIs, filtres, navettes\n• Cagnotte : contributions, messages, suivi"),
    
    ("2.2 Fonctionnalités Secondaires", 
     "• Communication : messagerie, notifications, forum\n• Gestion administrative : contrats, planning, budget, check-lists\n• Après événement : galerie, évaluations, remerciements"),
    
    ("3. Spécifications Techniques", 
     "Frontend : React, Tailwind, TypeScript\nBackend : Node.js, PostgreSQL, Stripe\nAPIs : Maps, Booking, Airbnb, Calendriers\nSécurité : HTTPS, RGPD, 2FA, logs"),
    
    ("4. Modules Fonctionnels", 
     "• Authentification : email, Google/Facebook, RGPD\n• Événements : CRUD, duplications, suivi\n• Professionnels : devis, réservation, comparaison\n• Invitations : création, envoi, relance, export\n• Hébergements : recherche, filtres, réservation\n• Cagnotte : création, suivi, contribution"),
    
    ("5. UX/UI & Design", 
     "• Design system cohérent\n• Mobile First et responsive\n• Couleurs : Rose, Bleu, Crème, Doré\n• Typo : Inter & Open Sans\n• Animations fluides, feedback utilisateur"),
    
    ("6. Contraintes", 
     "Techniques : navigateurs, SEO, performance\nLégales : RGPD, cookies, mentions légales\nSécurité : paiements, chiffrement, audit"),
    
    ("7. Planning & Livrables", 
     "Phase 1 (8 semaines) : MVP\nPhase 2 (6 semaines) : Avancées\nLivrables : code, app, doc technique, guide utilisateur, tests"),
    
    ("8. Budget & Ressources", 
     "Équipe : chef projet, devs, designer, QA\nCoûts : dev, APIs, hébergement, support"),
    
    ("9. Critères de Succès", 
     "• Techniques : performance, uptime, couverture de tests\n• Business : adoption, conversion, satisfaction\n• UX : temps de création < 5 min, support < 24h"),
    
    ("10. Risques & Mitigation", 
     "• Techniques : APIs, performance, sécurité\n• Business : concurrence, adoption, modèle\n• Projet : retards, budget, équipe"),
    
    ("11. Évolutions Futures", 
     "• V2 : app mobile, IA, AR, marketplace\n• Expansion : pro, international, partenariats")
]

for title, content in sections:
    add_slide(title, content)

# Sauvegarder le fichier
prs.save("EventMagic_Cahier_Des_Charges.pptx")
print("Fichier PowerPoint généré : EventMagic_Cahier_Des_Charges.pptx")
